#!/usr/bin/env python3
"""Stage-1 mechanical extraction: turn Office/legacy documents in raw/ into text.

The .docx/.pptx/.xlsx/.rtf/.doc sibling of scripts/extract-pdfs.sh and
scripts/extract-html.py. Each `foo.<ext>` gets a `foo.<ext>.txt` companion next to
it; the raw-backfill selector (scripts/cron/select_raw_batch.py) then prefers the
`.txt` over the binary so the ingest agent reads clean text. Domain-agnostic —
"doc in, text out".

  docx -> python-docx: body paragraphs + every table cell.
  pptx -> python-pptx: each slide's shape text (recursing groups) + tables + notes.
  xlsx -> openpyxl: non-empty cells per row, per sheet.
  rtf  -> striprtf: plain text.
  doc  -> antiword/catdoc (host tool; legacy binary has no pure-python reader).

Optional dependency, graceful skip (NO hard dep): each backend is imported / probed
per-format. If a lib/tool is absent that format is skipped (the others still run);
if NONE is available the script is a clean no-op (exit 0), so the host-side
extract-raw.sh wrapper never fails because a box lacks the office libs.

Idempotent: skips a source whose `.txt` companion is newer; re-extracts when the
source is touched. Empty output (image-only doc, the analogue of a scanned PDF
with no text layer) is a failure — no companion written, so the selector never
ingests an empty placeholder. Writes atomically (tmp + os.replace).

Deps (host):  python3 -m pip install python-docx python-pptx openpyxl striprtf
              + apt-get install antiword   (for legacy .doc)

Usage:
  python scripts/extract-docs.py                  # default: $WIKI_PATH/raw (else /opt/vault/raw)
  python scripts/extract-docs.py /path/to/raw
  python scripts/extract-docs.py --dry-run [root]
  python scripts/extract-docs.py --force [root]   # re-extract even if companion is newer
"""
from __future__ import annotations

import argparse
import os
import shutil
import subprocess
import sys
from pathlib import Path


def _build_backends() -> tuple[dict, list[str]]:
    """Return ({ext: extractor_fn}, [missing-lib names]) for whichever office
    libraries are installed. Importing a backend is optional — a missing lib just
    drops its format."""
    backends: dict = {}
    missing: list[str] = []

    try:
        import docx  # python-docx
    except ImportError:
        missing.append("python-docx (.docx)")
    else:
        def _docx_text(path: Path, _docx=docx) -> str:
            d = _docx.Document(str(path))
            parts: list[str] = [p.text for p in d.paragraphs]
            for table in d.tables:
                for row in table.rows:
                    parts.append(" | ".join(c.text for c in row.cells))
            return "\n".join(s for s in parts if s.strip())
        backends[".docx"] = _docx_text

    try:
        import pptx
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        missing.append("python-pptx (.pptx)")
    else:
        def _shape_text(shape) -> list[str]:
            """Recurse a pptx shape: grouped sub-shapes, text frames, tables."""
            out: list[str] = []
            try:
                is_group = shape.shape_type == MSO_SHAPE_TYPE.GROUP
            except Exception:
                is_group = False  # some shapes raise on .shape_type — treat as leaf
            if is_group:
                for sub in shape.shapes:
                    out.extend(_shape_text(sub))
                return out
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                out.append(shape.text_frame.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    out.append(" | ".join(c.text for c in row.cells))
            return out

        def _pptx_text(path: Path, _pptx=pptx) -> str:
            pr = _pptx.Presentation(str(path))
            parts: list[str] = []
            for slide in pr.slides:
                for shape in slide.shapes:
                    parts.extend(_shape_text(shape))
                if slide.has_notes_slide:
                    ntf = slide.notes_slide.notes_text_frame
                    if ntf is not None and ntf.text.strip():
                        parts.append("[notes] " + ntf.text)
            return "\n".join(s for s in parts if s.strip())
        backends[".pptx"] = _pptx_text

    try:
        import openpyxl
    except ImportError:
        missing.append("openpyxl (.xlsx)")
    else:
        def _xlsx_text(path: Path, _opx=openpyxl) -> str:
            wb = _opx.load_workbook(str(path), read_only=True, data_only=True)
            parts: list[str] = []
            for ws in wb.worksheets:
                rows = []
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None and str(c).strip()]
                    if cells:
                        rows.append(" | ".join(cells))
                if rows:
                    parts.append(f"# sheet: {ws.title}")
                    parts.extend(rows)
            wb.close()
            return "\n".join(parts)
        backends[".xlsx"] = _xlsx_text

    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        missing.append("striprtf (.rtf)")
    else:
        def _rtf_text(path: Path, _conv=rtf_to_text) -> str:
            return _conv(path.read_text(encoding="utf-8", errors="replace")).strip()
        backends[".rtf"] = _rtf_text

    # Legacy binary .doc has no pure-python reader — shell out to antiword/catdoc
    # (a host tool), if present.
    _doc_tool = shutil.which("antiword") or shutil.which("catdoc")
    if not _doc_tool:
        missing.append("antiword/catdoc (.doc)")
    else:
        def _doc_text(path: Path, _tool=_doc_tool) -> str:
            r = subprocess.run([_tool, str(path)], capture_output=True, timeout=60)
            return r.stdout.decode("utf-8", "replace").strip()
        backends[".doc"] = _doc_text

    return backends, missing


def main(argv: list[str]) -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("root", nargs="?", default="")
    ap.add_argument("--dry-run", action="store_true")
    ap.add_argument("--force", action="store_true", help="re-extract even if the companion is newer")
    args = ap.parse_args(argv)

    raw_root = Path(args.root or f"{os.environ.get('WIKI_PATH', '/opt/vault')}/raw")
    if not raw_root.is_dir():
        print(f"ERROR: raw root not found: {raw_root}", file=sys.stderr)
        return 1

    backends, missing = _build_backends()
    if not backends:
        # Neither office lib installed: clean no-op so the wrapper doesn't fail.
        print("note: no office extractors available — "
              "`pip install python-docx python-pptx` to extract .docx/.pptx. Skipping.")
        return 0
    if missing:
        print(f"note: {', '.join(missing)} not installed — those formats skipped this run.")

    exts = tuple(backends)
    extracted = skipped = failed = total = 0
    for src in sorted(raw_root.rglob("*")):
        if not src.is_file() or src.suffix.lower() not in exts:
            continue
        total += 1
        companion = src.with_name(src.name + ".txt")
        if not args.force and companion.is_file() and companion.stat().st_mtime >= src.stat().st_mtime:
            skipped += 1
            continue
        if args.dry_run:
            print(f"DRY: {src} -> {companion}")
            extracted += 1
            continue
        try:
            text = backends[src.suffix.lower()](src)
        except Exception as e:  # corrupt / unsupported variant — never abort the batch
            print(f"WARN: extract failed ({type(e).__name__}): {src}", file=sys.stderr)
            failed += 1
            continue
        if not text.strip():
            print(f"WARN: empty output (likely image-only, needs OCR): {src}", file=sys.stderr)
            failed += 1
            continue
        tmp = companion.with_name(companion.name + ".tmp")
        tmp.write_text(text + "\n", encoding="utf-8")
        os.replace(tmp, companion)            # atomic publish
        extracted += 1
        if extracted % 100 == 0:
            print(f"  ... {extracted} extracted, {skipped} skipped, {failed} failed")

    print(f"\nTotal: {total} office docs scanned ({'/'.join(e.lstrip('.') for e in exts)})")
    print(f"  extracted: {extracted}")
    print(f"  skipped (companion newer): {skipped}")
    print(f"  failed (no text / errored): {failed}")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
